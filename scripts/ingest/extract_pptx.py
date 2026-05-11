"""Extract text + speaker notes from BIOL 1592 unit PPTs into per-unit JSON.

Each input file is one unit (named like "08 Homeostasis - 2025.pptx").
Output: content/units/unit-XX.pptx.json — one JSON per unit, keyed by unit number
parsed from the filename prefix.

Designed to be re-run idempotently. Skips files that don't match the unit pattern.
"""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parents[2]
MATERIALS = ROOT / "materials"
OUT = ROOT / "content" / "units"

# "08 Homeostasis - 2025.pptx" -> (8, "Homeostasis")
# "00 Intro to BIOL 1592 - 2025.pptx" -> (0, "Intro to BIOL 1592")
# "07 Body Structures 2025.pptx" -> (7, "Body Structures")  (no dash)
FILENAME_RE = re.compile(r"^(\d{2})\s+(.*?)\s*(?:-\s*)?\d{4}\.pptx$", re.IGNORECASE)


def parse_unit_from_filename(name: str) -> tuple[int, str] | None:
    m = FILENAME_RE.match(name)
    if not m:
        return None
    return int(m.group(1)), m.group(2).strip()


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs)
        if line.strip():
            parts.append(line)
    return "\n".join(parts)


def extract_slide(slide, slide_no: int) -> dict:
    """Pull title, ordered body text, and speaker notes off a slide."""
    title = ""
    body_parts: list[str] = []
    image_alts: list[str] = []

    # python-pptx exposes .shapes.title for layouts that declare one;
    # not all slides do (some use a manual text box for the title).
    try:
        if slide.shapes.title is not None:
            t = shape_text(slide.shapes.title).strip()
            if t:
                title = t
    except (AttributeError, KeyError):
        pass

    # Sort shapes top-to-bottom, left-to-right so the extracted text follows
    # roughly the reading order. Shapes without position info sort last.
    def sort_key(s):
        top = s.top if s.top is not None else 10**9
        left = s.left if s.left is not None else 10**9
        return (top, left)

    for shape in sorted(slide.shapes, key=sort_key):
        # Skip the title shape (already captured) by identity.
        try:
            if shape is slide.shapes.title:
                continue
        except (AttributeError, KeyError):
            pass

        if shape.has_text_frame:
            text = shape_text(shape).strip()
            if text:
                body_parts.append(text)

        # Capture image alt text where present (for future drawing/identify support).
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                alt = shape.element.xpath(
                    ".//p:nvPicPr/p:cNvPr/@descr",
                    namespaces={"p": "http://schemas.openxmlformats.org/presentationml/2006/main"},
                )
                if alt and alt[0].strip():
                    image_alts.append(alt[0].strip())
            except Exception:
                pass

    notes = ""
    if slide.has_notes_slide:
        nt = slide.notes_slide.notes_text_frame.text if slide.notes_slide.notes_text_frame else ""
        notes = nt.strip()

    return {
        "slide_no": slide_no,
        "title": title,
        "body": "\n\n".join(body_parts),
        "notes": notes,
        "image_alts": image_alts,
    }


def extract_unit(pptx_path: Path, unit_id: int, unit_title: str) -> dict:
    prs = Presentation(str(pptx_path))
    slides = [extract_slide(slide, i + 1) for i, slide in enumerate(prs.slides)]
    return {
        "unit_id": unit_id,
        "title": unit_title,
        "source_file": str(pptx_path.relative_to(ROOT)),
        "slide_count": len(slides),
        "slides": slides,
    }


def main() -> int:
    if not MATERIALS.exists():
        print(f"materials/ not found at {MATERIALS}", file=sys.stderr)
        return 1

    OUT.mkdir(parents=True, exist_ok=True)

    pptx_files = sorted(MATERIALS.glob("*.pptx"))
    if not pptx_files:
        print("No .pptx files found in materials/", file=sys.stderr)
        return 1

    processed = 0
    for path in pptx_files:
        parsed = parse_unit_from_filename(path.name)
        if not parsed:
            print(f"  SKIP  unmatched filename: {path.name}")
            continue
        unit_id, unit_title = parsed
        out_path = OUT / f"unit-{unit_id:02d}.pptx.json"

        try:
            data = extract_unit(path, unit_id, unit_title)
        except Exception as e:
            print(f"  FAIL  unit {unit_id:02d} ({path.name}): {e}", file=sys.stderr)
            continue

        with out_path.open("w", encoding="utf-8") as fp:
            json.dump(data, fp, ensure_ascii=False, indent=2)

        body_chars = sum(len(s["body"]) for s in data["slides"])
        notes_chars = sum(len(s["notes"]) for s in data["slides"])
        print(
            f"  OK    unit {unit_id:02d}  {data['slide_count']:>3} slides  "
            f"body {body_chars:>6}c  notes {notes_chars:>5}c  → {out_path.name}"
        )
        processed += 1

    print(f"\nDone. {processed} unit(s) extracted to {OUT.relative_to(ROOT)}/")
    return 0


if __name__ == "__main__":
    sys.exit(main())
